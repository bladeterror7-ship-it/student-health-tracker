#!/usr/bin/env python3
"""WELLBE+ танилцуулга PPT үүсгэх: python3 scripts/generate-wellbe-ppt.py"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "WELLBE-taniltsuulga.pptx"
OUT_LEGACY = ROOT / "WELLBE+-taniltsuulga.pptx"

# Pastel / medical palette
EMERALD = RGBColor(0x05, 0x96, 0x69)
TEAL = RGBColor(0x0D, 0x94, 0x88)
SKY = RGBColor(0x02, 0x84, 0xC7)
VIOLET = RGBColor(0x6D, 0x28, 0xD9)
SLATE = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MINT_BG = RGBColor(0xEC, 0xFD, 0xF5)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MINT_BG
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = subtitle
    sub.font.size = Pt(20)
    sub.font.color.rgb = SLATE
    sub.alignment = PP_ALIGN.CENTER
    sub.space_before = Pt(12)


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    accent: RGBColor = EMERALD,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.9), Inches(0.9))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = accent

    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(8.7), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = SLATE
        p.level = 0
        p.space_after = Pt(10)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(30)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = TEAL

    for col, (ct, items, x) in enumerate(
        [(left_title, left_items, 0.5), (right_title, right_items, 5.1)]
    ):
        ht = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(4.3), Inches(0.5))
        ht.text_frame.paragraphs[0].text = ct
        ht.text_frame.paragraphs[0].font.size = Pt(20)
        ht.text_frame.paragraphs[0].font.bold = True
        ht.text_frame.paragraphs[0].font.color.rgb = EMERALD

        box = slide.shapes.add_textbox(Inches(x), Inches(1.75), Inches(4.3), Inches(4.8))
        tf = box.text_frame
        for i, line in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(16)
            p.font.color.rgb = SLATE
            p.space_after = Pt(8)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "WELLBE+",
        "Healthy Habits, Better Life\nСургуулийн эрүүл мэндийн цогц портал",
    )

    add_bullet_slide(
        prs,
        "Асуудал ба зорилго",
        [
            "Сургуулийн эрүүл мэндийн мэдээлэл тархсан, цаасан бүртгэлд үлддэг",
            "Сурагч, эмч, эцэг эх хоорондын холбоо сул, үзлэгийн түүх алга болох эрсдэлтэй",
            "WELLBE+ зорилго: нэг платформ дээр урьдчилан сэргийлэх, бүртгэх, хянах",
            "Үр дүн: илүү тодорхой эмчийн шийдвэр, сурагчийн идэвхтэй оролцоо",
        ],
    )

    add_bullet_slide(
        prs,
        "Системийн тойм",
        [
            "Вэб портал: wellbe-plus.com (локал: localhost:5173)",
            "Гурван үндсэн хэрэглэгч: Сурагч · Админ (эмч/багш) · Эцэг эх",
            "Дөрвөн модуль: Эмч · Сэтгэл зүй · Биеийн тамир · Vivera усны төсөл",
            "Пастель өнгө, бөөрөнхий UI — сургуулийн орчинд ээлтэй дизайн",
        ],
        SKY,
    )

    add_two_column_slide(
        prs,
        "Хэрэглэгчийн эрх",
        "Сурагч",
        [
            "Имэйлээр нэвтрэх, ангийн профайл",
            "Эмч: цаг захиалга, эмчээс асуух",
            "Эмчийн үзлэгийн түүх (шүд, нүд, амьсгал, даралт, пульс)",
            "Сэтгэл зүй: сэтгэл хөдлөл, тест, амьсгалын дасгал",
            "Биеийн тамир: идэвх, оноо, видео заавар",
        ],
        "Админ & Эцэг эх",
        [
            "Админ: бүх сурагчийн бүртгэл, үзлэг оруулах",
            "Интерактив эмчийн үзлэг (шүдийн зураг, хараа, ECG)",
            "Эмчийн асуултын хайрцаг, анхааруулга",
            "Эцэг эх: холбогдсон хүүхдийн тойм, Vivera хяналт",
        ],
    )

    add_bullet_slide(
        prs,
        "Эмчийн модуль — Сурагч",
        [
            "Эмчийн цаг захиалга (7 хоногийн хуваарь, цаг сонгох)",
            "Эмчээс асуух (нэр нууцлах сонголттой)",
            "Ерөнхий төлөв, сүүлийн үзлэгийн огноо",
            "Үзлэгийн түүх — олон удаагийн бүртгэл, өмнөх үзлэгүүд",
            "Эрүүл мэндийн анхааруулга харах",
        ],
        TEAL,
    )

    add_bullet_slide(
        prs,
        "Эмчийн модуль — Админ",
        [
            "Сурагч сонгоод «Шинэ үзлэг» — огноотой түүх үүсгэнэ",
            "Шүд: интерактив шүдний зураг (цоорсон, ломбодсон)",
            "Нүд: OD/OS хараа, Snellen хүснэгт, доогуур анхааруулга",
            "Цусны даралт: Systolic/Diastolic + механик хэмжүүр визуал",
            "Пульс: BPM + ECG монитор долгион (өндөр пульсэд хурдасна)",
            "Амьсгал: зүрх ууш, уушгин, ханиалга, сонсгол",
            "Тайлбарт нэгтгэх → сурагчид автоматаар харагдана",
        ],
        EMERALD,
    )

    add_bullet_slide(
        prs,
        "Сэтгэл зүйн модуль",
        [
            "Сэтгэл хөдлөлийн сан: 6 төрөл, тодорхойлолт, зохицуулах арга",
            "Хайрын хэлийн тест (30 асуулт, 5 хэл)",
            "Геометрийн сэтгэл зүйн тест (5 дүрс)",
            "Өдрийн сэтгэл, амьсгалын дасгал, нөөц материал",
            "Админ: сэтгэл зүйчийн самбар, захиалга, мэдэгдэл",
        ],
        VIOLET,
    )

    add_bullet_slide(
        prs,
        "Биеийн тамир & Vivera",
        [
            "Биеийн тамир: идэвх, оноо бүртгэл, админ хяналт",
            "Видео заавар upload (админ → сурагч үзнэ)",
            "Vivera усны төсөл: усны хэрэглээ, сурагч/эцэг эх хяналт",
            "Бүх модуль нэг WELLBE+ брэндийн доор нэгдсэн",
        ],
        SKY,
    )

    add_bullet_slide(
        prs,
        "Технологийн стек",
        [
            "Frontend: React 19, TypeScript, Vite, Tailwind CSS 4",
            "Анимаци: Framer Motion · График: Recharts",
            "Backend: Vercel Serverless API + Express (локал dev)",
            "Өгөгдлийн сан: Neon PostgreSQL (сурагч, портал, үзлэг, асуулт)",
            "Нэвтрэх: bcrypt, role-based (student / admin / parent)",
        ],
        SLATE,
    )

    add_bullet_slide(
        prs,
        "Өгөгдөл хадгалалт",
        [
            "clinical_exams: үзлэг бүр Neon дээр (огноо + JSON state)",
            "Сурагч бүрт олон үзлэгийн түүх — устгахгүй, хугацааны хязгааргүй",
            "Локал cache: хурдан UI, сервертэй автомат синк",
            "Эмчийн асуулт, бүртгэл, анхааруулга — төвлөрсөн удирдлага",
            "Production: GitHub → Vercel auto-deploy",
        ],
        TEAL,
    )

    add_bullet_slide(
        prs,
        "Аюулгүй байдал & хандалт",
        [
            "Эрхээр хязгаарласан хуудас (/dashboard, /admin, /parent)",
            "Админ бүртгэл: урилгын код (PE-ADMIN-2026)",
            "Эцэг эх: хүүхдийн бүртгэлтэй холбоотой данс",
            "Нууц үг: hash хадгалалт (Neon)",
            "Responsive вэб — компьютер, таблет, утас",
        ],
    )

    add_title_slide(
        prs,
        "Баярлалаа!",
        "WELLBE+ · wellbe-plus.com\nАсуулт байвал холбогдоорой",
    )

    prs.save(OUT)
    prs.save(OUT_LEGACY)
    print(f"✓ Хадгаллаа: {OUT}")
    print(f"✓ Хуучин нэр: {OUT_LEGACY}")


if __name__ == "__main__":
    main()
